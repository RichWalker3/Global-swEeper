#!/usr/bin/env python3
"""
Generate Global-swEep presentation with Global-e branding
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Global-e brand colors (approximated from typical branding)
GLOBALE_BLUE = RGBColor(0, 51, 102)      # Dark blue
GLOBALE_GREEN = RGBColor(0, 166, 81)     # Green accent
GLOBALE_GRAY = RGBColor(102, 102, 102)   # Text gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GLOBALE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = GLOBALE_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, callout=None):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = GLOBALE_BLUE
    header.line.fill.background()
    
    # Title on header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▪ {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = GLOBALE_GRAY
        p.space_after = Pt(12)
    
    # Callout box if provided
    if callout:
        callout_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(6.5), Inches(5.5), Inches(3), Inches(0.8)
        )
        callout_shape.fill.solid()
        callout_shape.fill.fore_color.rgb = GLOBALE_GREEN
        callout_shape.line.fill.background()
        
        callout_tf = callout_shape.text_frame
        callout_tf.paragraphs[0].text = callout
        callout_tf.paragraphs[0].font.size = Pt(16)
        callout_tf.paragraphs[0].font.bold = True
        callout_tf.paragraphs[0].font.color.rgb = WHITE
        callout_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        callout_tf.word_wrap = True
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = GLOBALE_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * num_rows)).table
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GLOBALE_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = GLOBALE_GRAY
    
    return slide

def add_thank_you_slide(prs):
    """Add final thank you slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Full blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = GLOBALE_BLUE
    bg.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions? Try it: npm run web"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "Global-swEep", "Automated Website Assessment Tool")
    
    # Slide 2: Problem → Solution
    add_content_slide(prs, "From 1 Hour to 10 Minutes", [
        "Website Assessments require 1+ hour of manual research",
        "Browse every page, hunt for third-party apps, check policies",
        "Easy to miss integrations buried in checkout or page source",
        "",
        "Global-swEep automates the research phase",
        "Crawls home, collections, PDPs, cart, checkout, policies",
        "Detects 30+ third-party integrations automatically",
        "Outputs structured evidence ready for Jira",
    ], callout="1 hour → 10 minutes")
    
    # Slide 3: How It Works
    add_content_slide(prs, "Simple 4-Step Workflow", [
        "Enter URL → Paste merchant website",
        "Run Assessment → Automated scrape (30-60 sec)",
        "Copy Prompt → Click 'Copy for Cursor'",
        "Get WA → Paste in Cursor, get formatted output",
    ], callout="npm run web")
    
    # Slide 4: What It Detects
    add_table_slide(prs, "Automatic Detection", 
        ["Category", "Examples"],
        [
            ["Platform", "Shopify, BigCommerce, Magento"],
            ["Red Flags", "Smile.io, Recharge, Competitors"],
            ["Payments", "Shop Pay, PayPal, Apple Pay, Google Pay"],
            ["BNPL", "Afterpay, Klarna, Affirm"],
            ["Subscriptions", "Recharge, Bold, Ordergroove"],
            ["Returns", "Loop, ReturnGO, Happy Returns"],
            ["Policies", "Return windows, warranties, final sale"],
        ]
    )
    
    # Slide 5: Validation
    add_table_slide(prs, "Tested Against Real WA Tickets",
        ["Metric", "Result"],
        [
            ["Merchants Tested", "24"],
            ["Successfully Scraped", "21 (87.5%)"],
            ["High Accuracy Matches", "5 sites (1-2 minor diffs)"],
            ["Platform Detection", "100%"],
            ["Third-Party Detection", "85%+"],
        ]
    )
    
    # Slide 6: Next Steps
    add_content_slide(prs, "Next Steps", [
        "Try It: npm run web → http://localhost:3847",
        "",
        "Pilot Proposal:",
        "  • Use for all new WA tickets this sprint",
        "  • Compare time spent vs. manual process",
        "  • Report edge cases for improvement",
        "",
        "Ask: Approve team pilot?",
    ])
    
    # Slide 7: Thank You
    add_thank_you_slide(prs)
    
    # Save
    output_path = os.path.expanduser("~/Desktop/Global-swEep_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
