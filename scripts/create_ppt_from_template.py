#!/usr/bin/env python3
"""
Generate Global-swEep presentation using Global-e template
Matches formatting of user's slide 2: Content Placeholder with levels
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def populate_content_placeholder(slide, content_lines):
    """
    Find the Content Placeholder and populate it with formatted content.
    Headers (ending with :) are level=0, bold, 20pt
    Regular items are level=1, 18pt
    Empty strings create blank paragraphs for spacing
    """
    content_ph = None
    for shape in slide.shapes:
        if "Content Placeholder" in shape.name:
            content_ph = shape
            break
    
    if not content_ph or not content_ph.has_text_frame:
        return
    
    tf = content_ph.text_frame
    tf.clear()
    
    first = True
    for line in content_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        if line == "":
            p.text = ""
        elif line.endswith(":"):
            # Header style
            p.text = line
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(20)
        else:
            # Bullet item style
            p.text = line
            p.level = 1
            p.font.size = Pt(18)

def set_title(slide, title_text):
    """Set the slide title"""
    for shape in slide.shapes:
        if "Title" in shape.name and shape.has_text_frame:
            shape.text_frame.paragraphs[0].text = title_text
            return

def main():
    template_path = "/Users/richard.walker/Downloads/Global-e GEM Integration Overview 2026.pptx"
    prs = Presentation(template_path)
    
    # Layouts
    title_layout = prs.slide_layouts[1]      # "Title Slide"
    content_layout = prs.slide_layouts[11]   # "Title and Content"
    thank_you_layout = prs.slide_layouts[10] # "Thank-you"
    
    # Clear existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(title_layout)
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.top.inches < 5:
                shape.text_frame.paragraphs[0].text = "Global-swEep"
            elif shape.top.inches < 9:
                shape.text_frame.paragraphs[0].text = "Automated Website Assessment Tool"
            elif shape.top.inches < 11:
                shape.text_frame.paragraphs[0].text = "March 2026"
    
    # ========== Slide 2: Problem → Solution ==========
    slide = prs.slides.add_slide(content_layout)
    set_title(slide, "From 1 Hour to 10 Minutes")
    populate_content_placeholder(slide, [
        "The Problem:",
        "Website Assessments require 1+ hour of manual research",
        "Browse every page, hunt for third-party apps, check policies",
        "Easy to miss integrations buried in checkout or page source",
        "",
        "The Solution:",
        "Global-swEep automates the research phase",
        "Crawls home, collections, PDPs, cart, checkout, policies",
        "Detects 30+ third-party integrations automatically",
        "Outputs structured evidence ready for Jira",
    ])
    
    # ========== Slide 3: How It Works ==========
    slide = prs.slides.add_slide(content_layout)
    set_title(slide, "Simple 4-Step Workflow")
    populate_content_placeholder(slide, [
        "Step 1:",
        "Enter URL → Paste merchant website",
        "",
        "Step 2:",
        "Run Assessment → Automated scrape (30-60 sec)",
        "",
        "Step 3:",
        "Copy Prompt → Click 'Copy for Cursor'",
        "",
        "Step 4:",
        "Get WA → Paste in Cursor, get formatted output",
        "",
        "No Setup Required:",
        "Run locally with: npm run web",
    ])
    
    # ========== Slide 4: What It Detects ==========
    slide = prs.slides.add_slide(content_layout)
    set_title(slide, "Automatic Detection")
    populate_content_placeholder(slide, [
        "Platform:",
        "Shopify, BigCommerce, Magento",
        "",
        "Red Flags:",
        "Smile.io, Recharge, Competitors",
        "",
        "Payments:",
        "Shop Pay, PayPal, Apple Pay, Google Pay",
        "",
        "BNPL:",
        "Afterpay, Klarna, Affirm",
        "",
        "Subscriptions:",
        "Recharge, Bold, Ordergroove",
        "",
        "Returns:",
        "Loop, ReturnGO, Happy Returns",
        "",
        "30+ apps detected automatically",
    ])
    
    # ========== Slide 5: Validation ==========
    slide = prs.slides.add_slide(content_layout)
    set_title(slide, "Tested Against Real WA Tickets")
    populate_content_placeholder(slide, [
        "Test Results:",
        "24 merchants tested",
        "7 excluded (URL failures, timeouts - not sweEp's fault)",
        "17 successfully scraped (100% of valid URLs)",
        "",
        "Accuracy:",
        "13/17 successful (76%) - matches or found MORE than Jira",
        "100% platform detection",
        "",
        "Key Finding:",
        "Scraper often finds MORE integrations than manual WAs",
    ])
    
    # ========== Slide 6: Next Steps ==========
    slide = prs.slides.add_slide(content_layout)
    set_title(slide, "Next Steps")
    populate_content_placeholder(slide, [
        "Try It:",
        "npm run web → http://localhost:3847",
        "",
        "Pilot Proposal:",
        "Use for all new WA tickets this sprint",
        "Compare time spent vs. manual process",
        "Report edge cases for improvement",
        "",
        "Ask:",
        "Approve team pilot?",
    ])
    
    # ========== Slide 7: What's Next ==========
    slide = prs.slides.add_slide(content_layout)
    set_title(slide, "What's Next")
    populate_content_placeholder(slide, [
        "Immediate Need:",
        "Global-e access to Cursor API / Anthropic API",
        "Removes copy-paste step → fully automated WA generation",
        "",
        "Planned Features:",
        "Direct LLM integration in web UI",
        "Jira API integration → auto-create WA tickets",
        "Batch processing for multiple merchants",
        "Enhanced third-party detection patterns",
        "",
        "Timeline:",
        "API access → Same-week integration",
        "Jira integration → Fast follow",
    ])
    
    # ========== Slide 8: Thank You ==========
    slide = prs.slides.add_slide(thank_you_layout)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(18), Inches(2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Questions? Try it: npm run web"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = os.path.expanduser("~/Desktop/Global-swEep_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
