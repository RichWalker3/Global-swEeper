#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches

template_path = "/Users/richard.walker/Downloads/Global-e GEM Integration Overview 2026.pptx"
prs = Presentation(template_path)

print(f"Template has {len(prs.slides)} slides\n")

# Look at first few slides to understand structure
for i, slide in enumerate(list(prs.slides)[:5]):
    print(f"=== Slide {i+1} ===")
    print(f"Layout: {slide.slide_layout.name}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text[:50].replace('\n', ' ') if shape.text_frame.text else "(empty)"
            print(f"  Shape: {shape.name}, pos=({shape.left.inches:.1f}, {shape.top.inches:.1f}), text: {text}")
    print()

# Also check slide dimensions
print(f"\nSlide dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
