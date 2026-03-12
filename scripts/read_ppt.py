#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt

ppt_path = "/Users/richard.walker/Desktop/Global-swEep_Presentation.pptx"
prs = Presentation(ppt_path)

for i, slide in enumerate(list(prs.slides)):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1} - Layout: {slide.slide_layout.name}")
    print('='*60)
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"\n  Shape: {shape.name}")
            print(f"  Position: ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
            print(f"  Size: ({shape.width.inches:.2f}\" x {shape.height.inches:.2f}\")")
            for j, para in enumerate(shape.text_frame.paragraphs):
                if para.text.strip():
                    font_size = para.font.size.pt if para.font.size else "inherit"
                    font_bold = para.font.bold if para.font.bold is not None else "inherit"
                    print(f"    Para {j}: \"{para.text[:60]}...\"" if len(para.text) > 60 else f"    Para {j}: \"{para.text}\"")
                    print(f"           size={font_size}, bold={font_bold}, level={para.level}")
